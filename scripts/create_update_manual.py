#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
QIP Dashboard V10 Admin 설정 관리 업데이트 매뉴얼 PPTX 생성 + 이메일 발송

Usage:
    # PPTX만 생성
    python scripts/create_update_manual.py --pptx-only

    # PPTX 생성 + 이메일 발송
    python scripts/create_update_manual.py --send
"""

import os
import sys
import json
import argparse
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from datetime import datetime, timezone

# Add project root to path for utils import
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
from scripts.utils.firebase_common import init_firestore

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("python-pptx 필요: pip install python-pptx")
    sys.exit(1)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "output_files")
PPTX_FILENAME = "QIP_Dashboard_V10_Admin_Manual.pptx"
SCREENSHOT_DIR = PROJECT_ROOT  # screenshots are in project root

SMTP_HOST = "mail.hsvina.com"
SMTP_PORT = 465

RECIPIENTS = ["hwk_qa@hsvina.com", "ksmoon@hsvina.com"]
DASHBOARD_URL = "https://moonkaicuzui.github.io/hwk-qip-incentive-v10/"

# Brand colors
COLOR_PRIMARY = RGBColor(0x35, 0x3A, 0x6E)  # Navy
COLOR_ACCENT = RGBColor(0x4A, 0x56, 0xC7)   # Indigo
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
COLOR_GREEN = RGBColor(0x28, 0xA7, 0x45)
COLOR_ORANGE = RGBColor(0xFD, 0x7E, 0x14)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def add_background(slide, color=COLOR_WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, y=Inches(0), height=Inches(1.2)):
    """Add a colored title bar at the top."""
    # Bar shape
    shape = slide.shapes.add_shape(
        1, Inches(0), y, SLIDE_WIDTH, height  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.15), Inches(11), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True


def add_subtitle(slide, text, y=Inches(1.2)):
    """Add subtitle text."""
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_GRAY


def add_bullet_list(slide, items, x=Inches(0.8), y=Inches(2.0), width=Inches(11), font_size=Pt(18)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(x, y, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (title, description)
            run = p.add_run()
            run.text = item[0]
            run.font.size = font_size
            run.font.bold = True
            run.font.color.rgb = COLOR_DARK

            run2 = p.add_run()
            run2.text = f"  —  {item[1]}"
            run2.font.size = Pt(font_size.pt - 2)
            run2.font.color.rgb = COLOR_GRAY
        else:
            p.text = item
            p.font.size = font_size
            p.font.color.rgb = COLOR_DARK

        p.space_before = Pt(8)
        p.space_after = Pt(4)
        p.level = 0


def add_screenshot(slide, img_path, x=Inches(0.5), y=Inches(1.8), max_width=Inches(12), max_height=Inches(5.2)):
    """Add a screenshot image, scaled to fit."""
    if not os.path.exists(img_path):
        # Add placeholder text
        txBox = slide.shapes.add_textbox(x, y, max_width, Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Screenshot: {os.path.basename(img_path)}]"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRAY
        return

    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size

    # Calculate scale to fit within bounds
    scale_w = max_width / Emu(img_w * 914400 / 96)  # px to EMU
    scale_h = max_height / Emu(img_h * 914400 / 96)
    scale = min(scale_w, scale_h, 1.0)

    width = Emu(int(img_w * 914400 / 96 * scale))
    height = Emu(int(img_h * 914400 / 96 * scale))

    slide.shapes.add_picture(img_path, x, y, width, height)


def add_numbered_steps(slide, steps, x=Inches(0.8), y=Inches(2.0), width=Inches(11)):
    """Add numbered step list."""
    txBox = slide.shapes.add_textbox(x, y, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run_num = p.add_run()
        run_num.text = f"{i+1}. "
        run_num.font.size = Pt(18)
        run_num.font.bold = True
        run_num.font.color.rgb = COLOR_ACCENT

        run_text = p.add_run()
        run_text.text = step
        run_text.font.size = Pt(18)
        run_text.font.color.rgb = COLOR_DARK

        p.space_before = Pt(12)
        p.space_after = Pt(4)


# ---------------------------------------------------------------------------
# PPTX Creation
# ---------------------------------------------------------------------------
def create_manual_pptx():
    """Create the update manual PPTX."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLOR_PRIMARY)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "HWK QIP Incentive Dashboard V10"
    p.font.size = Pt(44)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Admin 설정 관리 기능 업데이트 매뉴얼"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)

    p3 = tf2.add_paragraph()
    p3.text = "Admin Configuration Management Update Manual"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xEE)

    p4 = tf2.add_paragraph()
    p4.text = "Hướng dẫn cập nhật quản lý cấu hình Admin"
    p4.font.size = Pt(22)
    p4.font.color.rgb = RGBColor(0xAA, 0xAA, 0xEE)

    # Date & version
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p5 = tf3.paragraphs[0]
    p5.text = f"Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}  |  Version 10.0  |  HWK QA Team"
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(0x99, 0x99, 0xCC)

    # ========== Slide 2: Update Overview ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "업데이트 개요  |  Update Overview  |  Tổng quan cập nhật")

    items = [
        ("Admin 4탭 구조", "임계값 설정 / 설정 관리 / 시스템 / 데이터 조회"),
        ("TYPE-2 직급 매핑 관리", "TYPE-2 → TYPE-1 매핑 추가/수정/삭제 (Position Mapping CRUD)"),
        ("QIP Talent Pool 관리", "특별 보너스 멤버 추가/삭제, 설정 토글 (Talent Pool Member Management)"),
        ("Auditor 구역 매핑 관리", "AND/OR 조건 편집 모달로 구역 배정 (Area Assignment with Condition Editor)"),
        ("연속 근무월 조회", "읽기 전용 데이터 검색/조회 (Consecutive Months Read-only Lookup)"),
        ("변경 이력 추적", "모든 설정 변경 감사 기록 (Audit Trail for All Config Changes)"),
    ]
    add_bullet_list(slide, items, y=Inches(1.6))

    # ========== Slide 3: How to Access ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "접속 방법  |  How to Access  |  Cách truy cập")

    steps = [
        f"대시보드 접속: {DASHBOARD_URL}",
        "이메일/비밀번호로 로그인 (Login with email/password)",
        "우측 상단 [Admin] 버튼 클릭 (Click [Admin] button in top-right)",
        "4개 탭에서 필요한 기능 선택 (Select the desired tab)",
    ]
    add_numbered_steps(slide, steps, y=Inches(1.6))

    # Access note
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "※ Admin 권한: ksmoon@hsvina.com, huynhnhurgqa04@hsvina.com"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_ORANGE
    p.font.bold = True

    # ========== Slide 4: Tab Overview with Screenshot ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "탭 1: 임계값 설정  |  Thresholds  |  Ngưỡng")
    add_subtitle(slide, "인센티브 계산 기준값 및 근무일 수 재정의 관리")
    add_screenshot(slide, os.path.join(SCREENSHOT_DIR, "manual-tab1-thresholds.png"),
                   y=Inches(2.0), max_height=Inches(4.8))

    # ========== Slide 5: Config Management Tab ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "탭 2: 설정 관리  |  Config Management  |  Quản lý cấu hình")
    add_subtitle(slide, "TYPE-2 직급 매핑 + QIP Talent Pool + Auditor 구역 매핑")
    add_screenshot(slide, os.path.join(SCREENSHOT_DIR, "manual-tab2-configs.png"),
                   y=Inches(2.0), max_height=Inches(4.8))

    # ========== Slide 6: TYPE-2 Position Mapping Detail ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "TYPE-2 직급 매핑  |  Position Mapping  |  Ánh xạ chức vụ")

    items = [
        ("테이블 보기", "현재 등록된 TYPE-2 → TYPE-1 매핑 전체 목록 표시"),
        ("직급 추가", "하단 입력란에 직급명 입력 → 매핑 대상 선택 → [+ 추가] 클릭"),
        ("직급 삭제", "각 행 우측 휴지통 아이콘 클릭"),
        ("저장", "[저장] 버튼 클릭 → Firestore에 즉시 반영 + 변경 이력 기록"),
    ]
    add_bullet_list(slide, items, y=Inches(1.6), font_size=Pt(20))

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "매핑 대상 옵션: ASSEMBLY INSPECTOR, GROUP LEADER, LINE LEADER, (V) SUPERVISOR, A.MANAGER, MANAGER"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GRAY

    # ========== Slide 7: Talent Pool Detail ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "QIP Talent Pool  |  Talent Pool  |  Nhóm tài năng QIP")

    items = [
        ("설정 토글 (3개)", "자동 적용 / 기존 인센티브에 추가 / 조건 충족 필요"),
        ("멤버 테이블", "사번, 이름, 시작일, 종료일, 월 보너스(VND), 상태 표시"),
        ("멤버 추가", "하단 입력란에 사번/이름/보너스/종료일 입력 → [+ 추가] 클릭"),
        ("멤버 삭제", "각 행 우측 휴지통 아이콘 클릭"),
        ("저장", "[저장] 클릭 → Firestore 반영 + 변경 이력 기록"),
    ]
    add_bullet_list(slide, items, y=Inches(1.6), font_size=Pt(20))

    # ========== Slide 8: Auditor Area Mapping Detail ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "Auditor 구역 매핑  |  Area Mapping  |  Phân vùng Auditor")

    items = [
        ("Model Master", "전체 구역(ALL) 자동 적용 — 추가/삭제 가능"),
        ("Auditor 목록", "아코디언 형태로 각 Auditor의 담당 구역 조건 표시"),
        ("조건 편집", "아코디언 펼침 → [조건 편집] 클릭 → 모달 팝업"),
        ("AND/OR 조건", "조건 그룹(AND/OR/ALL) 추가 → 각 그룹 내 필터(BUILDING/REPACKING PO) 설정"),
        ("기본 불량률", "하단 기본 불량률 임계값(%) 설정"),
        ("저장", "[저장] 클릭 → Firestore 반영 + 변경 이력 기록"),
    ]
    add_bullet_list(slide, items, y=Inches(1.6), font_size=Pt(18))

    # ========== Slide 9: System & Data Lookup ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_title_bar(slide, "탭 3~4: 시스템 & 데이터 조회  |  System & Data Lookup")

    # Left side - System tab
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "시스템 탭 (System)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    items_sys = [
        "시스템 상태 모니터링",
        "이메일 보고서 수신자 관리",
        "변경 이력 (7컬럼: 유형 배지 추가)",
    ]
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items_sys):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(8)

    # Right side - Data Lookup tab
    txBox3 = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "데이터 조회 탭 (Data Lookup)"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_ACCENT

    items_data = [
        "연속 근무월 추적 (읽기 전용)",
        "이름/사번 검색 기능",
        "사번, 이름, 직급, 이전/현재 연속월, 인센티브 표시",
    ]
    txBox4 = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.5), Inches(2))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    for i, item in enumerate(items_data):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(8)

    # Screenshots side by side
    add_screenshot(slide, os.path.join(SCREENSHOT_DIR, "manual-tab3-system.png"),
                   x=Inches(0.3), y=Inches(4.0), max_width=Inches(6.2), max_height=Inches(3.2))
    add_screenshot(slide, os.path.join(SCREENSHOT_DIR, "manual-tab4-data.png"),
                   x=Inches(6.8), y=Inches(4.0), max_width=Inches(6.2), max_height=Inches(3.2))

    # ========== Slide 10: FAQ / Contact ==========
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLOR_PRIMARY)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "문의처  |  Contact  |  Liên hệ"
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True

    items = [
        ("Dashboard URL", DASHBOARD_URL),
        ("Admin 문의", "ksmoon@hsvina.com"),
        ("QA Team", "hwk_qa@hsvina.com"),
        ("지원 언어", "한국어 (KO) / English (EN) / Tiếng Việt (VN)"),
    ]
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, (label, value) in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = f"{label}:  "
        run1.font.size = Pt(20)
        run1.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(20)
        run2.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(16)

    # Save
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, PPTX_FILENAME)
    prs.save(output_path)
    print(f"✅ PPTX 생성 완료: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Email HTML
# ---------------------------------------------------------------------------
def _email_wrapper(header_subtitle, body_content, footer_text):
    """Wrap email body content with header and footer."""
    return f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width"></head>
<body style="font-family:'Segoe UI',Arial,sans-serif;margin:0;padding:0;background:#f5f5f5;">
<table width="100%" cellpadding="0" cellspacing="0" style="max-width:700px;margin:0 auto;background:#ffffff;">
<tr><td style="background:#353a6e;padding:30px 40px;">
<h1 style="color:#ffffff;margin:0;font-size:24px;">HWK QIP Incentive Dashboard V10</h1>
<p style="color:#ccccff;margin:8px 0 0;font-size:14px;">{header_subtitle}</p>
</td></tr>
{body_content}
<tr><td style="background:#f8f9fa;padding:20px 40px;text-align:center;">
<p style="color:#999;font-size:12px;margin:0;">{footer_text}</p>
</td></tr>
</table>
</body>
</html>"""


def generate_email_html_ko():
    """한국어 전용 이메일 HTML."""
    body = f"""<tr><td style="padding:30px 40px;">
<h2 style="color:#353a6e;font-size:20px;margin:0 0 15px;">업데이트 안내</h2>
<p style="color:#333;line-height:1.6;">안녕하세요,<br>QIP Dashboard V10 Admin 페이지에 <b>설정 관리 기능</b>이 추가되었습니다.</p>
<h3 style="color:#4a56c7;font-size:16px;">주요 변경 사항</h3>
<ul style="color:#333;line-height:1.8;">
<li><b>Admin 4탭 구조</b> — 임계값 설정 / 설정 관리 / 시스템 / 데이터 조회</li>
<li><b>TYPE-2 직급 매핑</b> — TYPE-2 → TYPE-1 매핑 추가/수정/삭제</li>
<li><b>QIP Talent Pool</b> — 특별 보너스 멤버 관리 (추가/삭제/설정)</li>
<li><b>Auditor 구역 매핑</b> — AND/OR 조건 편집 모달로 구역 배정 관리</li>
<li><b>연속 근무월 조회</b> — 읽기 전용 데이터 검색/조회</li>
<li><b>변경 이력 추적</b> — 모든 설정 변경에 대한 감사 기록</li>
</ul>
<p style="color:#333;">📎 첨부된 PPTX 매뉴얼을 참고해 주세요.</p>
<p><a href="{DASHBOARD_URL}" style="display:inline-block;background:#4a56c7;color:#fff;padding:10px 24px;border-radius:6px;text-decoration:none;font-weight:bold;">대시보드 접속</a></p>
</td></tr>"""
    return _email_wrapper(
        "Admin 설정 관리 기능 업데이트", body,
        f"HWK QIP Incentive System V10 — {datetime.now(timezone.utc).strftime('%Y-%m-%d')}<br>문의: ksmoon@hsvina.com"
    )


def generate_email_html_en():
    """영어 전용 이메일 HTML."""
    body = f"""<tr><td style="padding:30px 40px;">
<h2 style="color:#353a6e;font-size:20px;margin:0 0 15px;">Update Notice</h2>
<p style="color:#333;line-height:1.6;">Hello,<br>The <b>Configuration Management</b> feature has been added to the QIP Dashboard V10 Admin page.</p>
<h3 style="color:#4a56c7;font-size:16px;">Key Changes</h3>
<ul style="color:#333;line-height:1.8;">
<li><b>Admin 4-Tab Layout</b> — Thresholds / Config Management / System / Data Lookup</li>
<li><b>TYPE-2 Position Mapping</b> — Add/Edit/Delete TYPE-2 → TYPE-1 mappings</li>
<li><b>QIP Talent Pool</b> — Manage bonus members (add/remove/settings)</li>
<li><b>Auditor Area Mapping</b> — AND/OR condition editor modal for area assignments</li>
<li><b>Consecutive Months Lookup</b> — Read-only data search and view</li>
<li><b>Change History</b> — Audit trail for all configuration changes</li>
</ul>
<p style="color:#333;">📎 Please refer to the attached PPTX manual.</p>
<p><a href="{DASHBOARD_URL}" style="display:inline-block;background:#4a56c7;color:#fff;padding:10px 24px;border-radius:6px;text-decoration:none;font-weight:bold;">Open Dashboard</a></p>
</td></tr>"""
    return _email_wrapper(
        "Admin Configuration Management Update", body,
        f"HWK QIP Incentive System V10 — {datetime.now(timezone.utc).strftime('%Y-%m-%d')}<br>Contact: ksmoon@hsvina.com"
    )


def generate_email_html_vn():
    """베트남어 전용 이메일 HTML."""
    body = f"""<tr><td style="padding:30px 40px;">
<h2 style="color:#353a6e;font-size:20px;margin:0 0 15px;">Thông báo cập nhật</h2>
<p style="color:#333;line-height:1.6;">Xin chào,<br>Tính năng <b>Quản lý cấu hình</b> đã được thêm vào trang Admin của QIP Dashboard V10.</p>
<h3 style="color:#4a56c7;font-size:16px;">Các thay đổi chính</h3>
<ul style="color:#333;line-height:1.8;">
<li><b>Admin 4 tab</b> — Ngưỡng / Quản lý cấu hình / Hệ thống / Tra cứu dữ liệu</li>
<li><b>Ánh xạ chức vụ TYPE-2</b> — Thêm/Sửa/Xóa ánh xạ TYPE-2 → TYPE-1</li>
<li><b>QIP Talent Pool</b> — Quản lý thành viên thưởng (thêm/xóa/cài đặt)</li>
<li><b>Phân vùng Auditor</b> — Trình chỉnh sửa điều kiện AND/OR cho phân vùng</li>
<li><b>Tra cứu tháng liên tục</b> — Tìm kiếm và xem dữ liệu (chỉ đọc)</li>
<li><b>Lịch sử thay đổi</b> — Nhật ký kiểm toán cho tất cả thay đổi cấu hình</li>
</ul>
<p style="color:#333;">📎 Vui lòng tham khảo tài liệu PPTX đính kèm.</p>
<p><a href="{DASHBOARD_URL}" style="display:inline-block;background:#4a56c7;color:#fff;padding:10px 24px;border-radius:6px;text-decoration:none;font-weight:bold;">Mở Dashboard</a></p>
</td></tr>"""
    return _email_wrapper(
        "Cập nhật quản lý cấu hình Admin", body,
        f"HWK QIP Incentive System V10 — {datetime.now(timezone.utc).strftime('%Y-%m-%d')}<br>Liên hệ: ksmoon@hsvina.com"
    )


# ---------------------------------------------------------------------------
# Email Sending
# ---------------------------------------------------------------------------
def _smtp_connect(smtp_user, smtp_password):
    """한비로 SMTP AUTH LOGIN 직접 연결 (CRAM-MD5/PLAIN 실패 방지)."""
    import ssl
    import base64
    ctx = ssl.create_default_context()
    server = smtplib.SMTP_SSL(SMTP_HOST, SMTP_PORT, timeout=30, context=ctx)
    server.ehlo()
    code, resp = server.docmd(
        "AUTH LOGIN",
        base64.b64encode(smtp_user.encode()).decode()
    )
    if code == 334:
        code, resp = server.docmd(
            base64.b64encode(smtp_password.encode()).decode()
        )
    if code != 235:
        raise Exception(f"AUTH LOGIN 실패: {code} {resp}")
    return server


def _send_one_email(smtp_user, smtp_password, recipient, subject, html_body, pptx_path):
    """단일 이메일 발송 (개별 SMTP 연결)."""
    server = _smtp_connect(smtp_user, smtp_password)

    msg = MIMEMultipart()
    msg["Subject"] = subject
    msg["From"] = f"QIP Dashboard <{smtp_user}>"
    msg["To"] = recipient
    msg.attach(MIMEText(html_body, "html", "utf-8"))

    # Attach PPTX
    with open(pptx_path, "rb") as f:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(f.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename={PPTX_FILENAME}")
        msg.attach(part)

    server.mail(smtp_user)
    server.rcpt(recipient)
    server.data(msg.as_bytes())
    server.quit()


def send_update_email(pptx_path):
    """언어별 3통 × 수신자 2명 = 총 6통 이메일 발송."""
    smtp_user = os.environ.get("SMTP_USER", "")
    smtp_password = os.environ.get("SMTP_PASSWORD", "")

    if not smtp_user or not smtp_password:
        # Try Firestore
        try:
            from firebase_admin import firestore
            db = init_firestore()
            config = db.collection("system").document("config").get()
            if config.exists:
                data = config.to_dict()
                email_settings = data.get("email_settings", {})
                smtp_user = smtp_user or email_settings.get("smtp_user", "")
                smtp_password = smtp_password or email_settings.get("smtp_password", "")
        except Exception as e:
            print(f"  Firestore 설정 로드 실패: {e}")

    if not smtp_user or not smtp_password:
        print("❌ SMTP 인증 정보 없음!")
        print("  SMTP_USER, SMTP_PASSWORD 환경변수를 설정하세요.")
        return False

    # 언어별 이메일 정의
    emails = [
        {
            "lang": "KO",
            "subject": "[QIP Dashboard V10] Admin 설정 관리 기능 업데이트 안내",
            "html_fn": generate_email_html_ko,
        },
        {
            "lang": "EN",
            "subject": "[QIP Dashboard V10] Admin Config Management Update",
            "html_fn": generate_email_html_en,
        },
        {
            "lang": "VN",
            "subject": "[QIP Dashboard V10] Cập nhật quản lý cấu hình Admin",
            "html_fn": generate_email_html_vn,
        },
    ]

    # Save HTML previews
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    for e in emails:
        html = e["html_fn"]()
        preview_path = os.path.join(OUTPUT_DIR, f"update_email_{e['lang'].lower()}.html")
        with open(preview_path, "w", encoding="utf-8") as f:
            f.write(html)
        print(f"  HTML 미리보기: {preview_path}")

    result = {"sent": 0, "failed": 0}

    for email_def in emails:
        lang = email_def["lang"]
        subject = email_def["subject"]
        html_body = email_def["html_fn"]()

        for recipient in RECIPIENTS:
            try:
                print(f"  [{lang}] {SMTP_HOST}:{SMTP_PORT} → {recipient}")
                _send_one_email(smtp_user, smtp_password, recipient, subject, html_body, pptx_path)
                result["sent"] += 1
                print(f"    ✅ {recipient}")
            except Exception as e:
                result["failed"] += 1
                print(f"    ❌ {recipient}: {e}")

    total = len(emails) * len(RECIPIENTS)
    print(f"\n  발송 결과: 성공 {result['sent']}/{total}, 실패 {result['failed']}/{total}")
    return result["sent"] > 0


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description="V10 Admin 업데이트 매뉴얼 + 이메일 발송")
    parser.add_argument("--pptx-only", action="store_true", help="PPTX만 생성 (이메일 미발송)")
    parser.add_argument("--send", action="store_true", help="PPTX 생성 + 이메일 발송")
    args = parser.parse_args()

    print("=" * 60)
    print("  QIP Dashboard V10 - Admin 업데이트 매뉴얼")
    print(f"  시간: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)

    # Step 1: Create PPTX
    print("\n📄 PPTX 매뉴얼 생성 중...")
    pptx_path = create_manual_pptx()

    if args.pptx_only:
        print("\n✅ 완료 (PPTX만 생성)")
        return

    if args.send:
        # Step 2: Send email
        print(f"\n📧 이메일 발송 중... → {', '.join(RECIPIENTS)}")
        send_update_email(pptx_path)
    else:
        print(f"\n  이메일 발송하려면: python scripts/create_update_manual.py --send")
        print(f"  SMTP_USER, SMTP_PASSWORD 환경변수 필요")

    print("\n✅ 완료")


if __name__ == "__main__":
    main()
