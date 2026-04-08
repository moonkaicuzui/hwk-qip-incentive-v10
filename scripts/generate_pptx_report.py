#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
February 2026 QIP Incentive Report — PPTX Generator for 전자결재
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           "output_files", "2026년_2월_QIP_인센티브_보고서.pptx")

# ── Colors ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)

# ── Data ──
SUMMARY = {
    "month": "2월", "year": "2026",
    "total_employees": 429, "eligible": 414, "receiving": 356,
    "not_receiving": 58, "excluded": 15,
    "total_incentive": 187_124_241,
    "working_days": 19,
    "pass_rate": 85.99,  # 356/414
    "data_updated": "2026-03-10",
}

TYPES = [
    {"name": "TYPE-1 (관리직)", "total": 135, "eligible": 135, "receiving": 109,
     "not_receiving": 26, "amount": 67_020_398, "pass_rate": 80.74},
    {"name": "TYPE-2 (생산직)", "total": 279, "eligible": 279, "receiving": 247,
     "not_receiving": 32, "amount": 120_103_843, "pass_rate": 88.53},
    {"name": "TYPE-3 (수습/면제)", "total": 15, "eligible": 0, "receiving": 0,
     "not_receiving": 0, "amount": 0, "pass_rate": 0},
]

BUILDINGS = [
    ("A1A", 50, 45, 23_038_397), ("A1B", 28, 25, 14_290_040),
    ("A2", 52, 42, 19_975_508), ("A3", 24, 22, 9_850_754),
    ("A4", 22, 17, 7_810_398), ("A5A", 9, 7, 3_780_214),
    ("A5B", 5, 4, 1_876_321), ("A6A", 2, 2, 1_120_000),
    ("A6B", 4, 3, 1_512_608), ("A7", 3, 3, 1_640_000),
    ("B1", 39, 34, 16_780_195), ("B2A", 21, 19, 9_180_321),
    ("B2B", 8, 7, 3_210_450), ("B3", 28, 24, 12_280_194),
    ("B5", 24, 22, 10_980_321), ("C1A", 25, 22, 11_220_408),
    ("C1B", 18, 16, 8_120_195), ("OFFICE", 67, 42, 30_457_917),
]

CONDITIONS = [
    ("1. 출근율 (≥90%)", 399, 15, 15, 96.4),
    ("2. AQL 검사 합격률 (≤2%)", 258, 21, 150, 92.5),
    ("3. 5PRS 통과율 (≥97%)", 252, 27, 150, 90.3),
    ("4. 연속 불합격 (≤2회)", 271, 8, 150, 97.1),
    ("5. 근속 기간 (≥6개월)", 407, 7, 15, 98.3),
    ("6. 징계 이력 (0건)", 412, 2, 15, 99.5),
    ("7. 무단 결근 (≤1일)", 409, 5, 15, 98.8),
    ("8. 지각 횟수 (≤3회)", 404, 10, 15, 97.6),
    ("9. 조기 퇴근 (≤2회)", 408, 6, 15, 98.6),
    ("10. 특별 감점 (0건)", 411, 3, 15, 99.3),
]

THRESHOLDS = [
    ("출근율", "≥ 90%"),
    ("AQL 불합격률", "≤ 2.0%"),
    ("5PRS 통과율", "≥ 97%"),
    ("5PRS 최소 검사수량", "≥ 150건"),
    ("무단 결근", "≤ 1일"),
    ("최소 근무일수", "≥ 12일"),
]


def fmt_krw(n):
    """Format number with commas + VND suffix."""
    return f"{n:,.0f} VND"


def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Add a colored rectangle as background."""
    from pptx.util import Inches as I
    w = width or I(13.33)
    h = height or I(7.5)
    shape = slide.shapes.add_shape(1, left, top, w, h)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def set_cell_text(cell, text, font_size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    """Set table cell text with formatting."""
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "맑은 고딕"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_bg(cell, color):
    """Set cell background color."""
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', '%02x%02x%02x' % (color[0] if isinstance(color, tuple) else color.red,
                                            color[1] if isinstance(color, tuple) else color.green,
                                            color[2] if isinstance(color, tuple) else color.blue))


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_kpi_card(slide, left, top, width, height, label, value, sub_text=None, accent_color=ACCENT_BLUE):
    """Add a KPI card with label and value."""
    from pptx.oxml.ns import qn
    # Card background
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)
    # Round corners
    shape.shadow.inherit = False

    # Accent bar on top
    bar = slide.shapes.add_shape(1, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.3),
                 label, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Value
    add_text_box(slide, left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), Inches(0.5),
                 str(value), font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    if sub_text:
        add_text_box(slide, left + Inches(0.1), top + Inches(0.85), width - Inches(0.2), Inches(0.25),
                     sub_text, font_size=8, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def create_report():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════
    # Slide 1: Title
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg_rect(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "QIP INCENTIVE 보고서", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "2026년 2월 — Quality Incentive Program 실적 분석",
                 font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(1, Inches(4), Inches(3.8), Inches(5.33), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
                 f"기준일: {SUMMARY['data_updated']}  |  근무일수: {SUMMARY['working_days']}일  |  대상: {SUMMARY['total_employees']}명",
                 font_size=14, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "HWK 품질관리팀", font_size=14, color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # Slide 2: 핵심 요약
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "핵심 요약 (Executive Summary)", font_size=24, bold=True, color=NAVY)

    # Divider
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # KPI Cards - Row 1
    card_w = Inches(2.5)
    card_h = Inches(1.2)
    y1 = Inches(1.3)
    gap = Inches(0.3)
    x_start = Inches(0.5)

    add_kpi_card(slide, x_start, y1, card_w, card_h,
                 "전체 인원", f"{SUMMARY['total_employees']}명",
                 f"대상: {SUMMARY['eligible']}명 / 제외: {SUMMARY['excluded']}명")
    add_kpi_card(slide, x_start + card_w + gap, y1, card_w, card_h,
                 "수령 인원", f"{SUMMARY['receiving']}명",
                 f"미수령: {SUMMARY['not_receiving']}명", GREEN)
    add_kpi_card(slide, x_start + (card_w + gap) * 2, y1, card_w, card_h,
                 "조건 통과율", f"{SUMMARY['pass_rate']:.1f}%",
                 f"{SUMMARY['receiving']}/{SUMMARY['eligible']}", ACCENT_BLUE)
    add_kpi_card(slide, x_start + (card_w + gap) * 3, y1, card_w, card_h,
                 "총 인센티브 금액", fmt_krw(SUMMARY['total_incentive']),
                 f"근무일수: {SUMMARY['working_days']}일", ORANGE)

    # TYPE breakdown table
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(5), Inches(0.4),
                 "TYPE별 현황", font_size=16, bold=True, color=NAVY)

    rows, cols = 4, 6
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.4), Inches(7.5), Inches(1.6))
    tbl = tbl_shape.table

    # Set column widths
    widths = [Inches(2.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.2), Inches(1.3)]
    for i, w in enumerate(widths):
        tbl.columns[i].width = w

    # Header
    headers = ["구분", "전체", "수령", "미수령", "통과율", "금액"]
    for i, h in enumerate(headers):
        set_cell_text(tbl.cell(0, i), h, font_size=10, bold=True, color=WHITE)
        set_cell_bg(tbl.cell(0, i), NAVY)

    for r, t in enumerate(TYPES, 1):
        set_cell_text(tbl.cell(r, 0), t["name"], font_size=10, alignment=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(r, 1), str(t["total"]), font_size=10)
        set_cell_text(tbl.cell(r, 2), str(t["receiving"]), font_size=10, color=GREEN)
        set_cell_text(tbl.cell(r, 3), str(t["not_receiving"]), font_size=10,
                      color=RED if t["not_receiving"] > 0 else DARK_GRAY)
        set_cell_text(tbl.cell(r, 4), f"{t['pass_rate']:.1f}%", font_size=10)
        set_cell_text(tbl.cell(r, 5), fmt_krw(t["amount"]), font_size=9, alignment=PP_ALIGN.RIGHT)
        # Alternating row color
        if r % 2 == 0:
            for c in range(cols):
                set_cell_bg(tbl.cell(r, c), LIGHT_GRAY)

    # Summary box on the right
    add_text_box(slide, Inches(8.5), Inches(2.9), Inches(4.3), Inches(0.4),
                 "주요 지표", font_size=16, bold=True, color=NAVY)

    summary_lines = [
        f"• 전체 통과율: {SUMMARY['pass_rate']:.1f}% ({SUMMARY['receiving']}/{SUMMARY['eligible']}명)",
        f"• TYPE-1 (관리직) 통과율: 80.7% — 26명 미수령",
        f"• TYPE-2 (생산직) 통과율: 88.5% — 32명 미수령",
        f"• TYPE-3 (수습/면제): 15명 전원 대상 제외",
        f"",
        f"• 총 인센티브: {fmt_krw(SUMMARY['total_incentive'])}",
        f"  - TYPE-1: {fmt_krw(67_020_398)} (35.8%)",
        f"  - TYPE-2: {fmt_krw(120_103_843)} (64.2%)",
    ]
    for i, line in enumerate(summary_lines):
        if line:
            add_text_box(slide, Inches(8.5), Inches(3.4) + Inches(i * 0.28), Inches(4.3), Inches(0.3),
                         line, font_size=10, color=DARK_GRAY)

    # ═══════════════════════════════════════════════════════
    # Slide 3: 10개 조건별 통과율
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "인센티브 조건별 통과율 분석", font_size=24, bold=True, color=NAVY)

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Conditions table
    rows, cols = 11, 6
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(8.5), Inches(5.5))
    tbl = tbl_shape.table

    widths = [Inches(2.8), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.2), Inches(1.5)]
    for i, w in enumerate(widths):
        tbl.columns[i].width = w

    headers = ["조건", "통과", "불합격", "해당없음", "통과율", "비고"]
    for i, h in enumerate(headers):
        set_cell_text(tbl.cell(0, i), h, font_size=10, bold=True, color=WHITE)
        set_cell_bg(tbl.cell(0, i), NAVY)

    for r, (name, passed, failed, na, rate) in enumerate(CONDITIONS, 1):
        set_cell_text(tbl.cell(r, 0), name, font_size=9, alignment=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(r, 1), str(passed), font_size=10, color=GREEN)
        set_cell_text(tbl.cell(r, 2), str(failed), font_size=10, color=RED if failed > 0 else DARK_GRAY)
        set_cell_text(tbl.cell(r, 3), str(na), font_size=10, color=MID_GRAY)
        set_cell_text(tbl.cell(r, 4), f"{rate:.1f}%", font_size=10, bold=True,
                      color=GREEN if rate >= 95 else (ORANGE if rate >= 90 else RED))
        # Remark
        remark = ""
        if rate < 93:
            remark = "주의 필요"
        elif rate >= 99:
            remark = "우수"
        set_cell_text(tbl.cell(r, 5), remark, font_size=9,
                      color=RED if "주의" in remark else GREEN)
        if r % 2 == 0:
            for c in range(cols):
                set_cell_bg(tbl.cell(r, c), LIGHT_GRAY)

    # Insight box
    add_text_box(slide, Inches(9.3), Inches(1.3), Inches(3.5), Inches(0.4),
                 "분석 인사이트", font_size=14, bold=True, color=NAVY)

    insights = [
        "가장 낮은 통과율:",
        "  3. 5PRS 통과율: 90.3%",
        "  2. AQL 합격률: 92.5%",
        "",
        "가장 높은 통과율:",
        "  6. 징계 이력: 99.5%",
        "  10. 특별 감점: 99.3%",
        "",
        "개선 포인트:",
        "  • 5PRS/AQL — 품질 교육 강화",
        "  • 27명 5PRS 불합격 집중 관리",
        "  • 21명 AQL 불합격 원인 분석",
    ]
    for i, txt in enumerate(insights):
        if txt:
            c = RED if "낮은" in txt else (GREEN if "높은" in txt else DARK_GRAY)
            add_text_box(slide, Inches(9.3), Inches(1.8) + Inches(i * 0.28), Inches(3.5), Inches(0.3),
                         txt, font_size=9, color=c)

    # ═══════════════════════════════════════════════════════
    # Slide 4: 건물별 현황
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "건물별 인센티브 현황", font_size=24, bold=True, color=NAVY)

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Split buildings into two tables (9 each)
    half = 9
    for tbl_idx in range(2):
        x_offset = Inches(0.5) + tbl_idx * Inches(6.3)
        bldg_subset = BUILDINGS[tbl_idx * half: (tbl_idx + 1) * half]
        r_count = len(bldg_subset) + 1

        tbl_shape = slide.shapes.add_table(r_count, 5, x_offset, Inches(1.3), Inches(6.0), Inches(0.35 * r_count))
        tbl = tbl_shape.table

        widths = [Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(2.0)]
        for i, w in enumerate(widths):
            tbl.columns[i].width = w

        headers = ["건물", "대상", "수령", "통과율", "인센티브 금액"]
        for i, h in enumerate(headers):
            set_cell_text(tbl.cell(0, i), h, font_size=9, bold=True, color=WHITE)
            set_cell_bg(tbl.cell(0, i), NAVY)

        for r, (bldg, total, recv, amt) in enumerate(bldg_subset, 1):
            rate = (recv / total * 100) if total > 0 else 0
            set_cell_text(tbl.cell(r, 0), bldg, font_size=9, bold=True)
            set_cell_text(tbl.cell(r, 1), str(total), font_size=9)
            set_cell_text(tbl.cell(r, 2), str(recv), font_size=9, color=GREEN)
            set_cell_text(tbl.cell(r, 3), f"{rate:.0f}%", font_size=9,
                          color=GREEN if rate >= 90 else (ORANGE if rate >= 80 else RED))
            set_cell_text(tbl.cell(r, 4), fmt_krw(amt), font_size=8, alignment=PP_ALIGN.RIGHT)
            if r % 2 == 0:
                for c in range(5):
                    set_cell_bg(tbl.cell(r, c), LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════
    # Slide 5: 적용 기준 & 결론
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "적용 기준 및 결론", font_size=24, bold=True, color=NAVY)

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Thresholds table
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(5), Inches(0.4),
                 "2월 적용 임계값", font_size=16, bold=True, color=NAVY)

    tbl_shape = slide.shapes.add_table(len(THRESHOLDS) + 1, 2, Inches(0.5), Inches(1.7), Inches(5), Inches(2.3))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(2.5)

    set_cell_text(tbl.cell(0, 0), "항목", font_size=10, bold=True, color=WHITE)
    set_cell_text(tbl.cell(0, 1), "기준값", font_size=10, bold=True, color=WHITE)
    set_cell_bg(tbl.cell(0, 0), NAVY)
    set_cell_bg(tbl.cell(0, 1), NAVY)

    for r, (item, val) in enumerate(THRESHOLDS, 1):
        set_cell_text(tbl.cell(r, 0), item, font_size=10, alignment=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(r, 1), val, font_size=10, bold=True, color=ACCENT_BLUE)
        if r % 2 == 0:
            set_cell_bg(tbl.cell(r, 0), LIGHT_GRAY)
            set_cell_bg(tbl.cell(r, 1), LIGHT_GRAY)

    # Conclusion box
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(0.4),
                 "결론 및 건의사항", font_size=16, bold=True, color=NAVY)

    # Conclusion card background
    card = slide.shapes.add_shape(1, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BLUE
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(1)

    conclusions = [
        ("1. 인센티브 지급 승인 요청", True),
        (f"   총 {SUMMARY['receiving']}명, {fmt_krw(SUMMARY['total_incentive'])} 지급 건의", False),
        ("", False),
        ("2. 전체 통과율 85.99%", True),
        (f"   414명 중 356명 조건 충족, 58명 미충족", False),
        ("", False),
        ("3. 품질 개선 포인트", True),
        ("   • 5PRS 통과율 (90.3%) — 27명 불합격", False),
        ("   • AQL 합격률 (92.5%) — 21명 불합격", False),
        ("   → 품질 교육 강화 및 불합격자 집중 관리 필요", False),
        ("", False),
        ("4. TYPE별 현황", True),
        ("   • 관리직(TYPE-1): 80.7% — 상대적 저조", False),
        ("   • 생산직(TYPE-2): 88.5% — 양호", False),
        ("", False),
        ("5. 특이사항", True),
        ("   • OFFICE 건물: 67명 중 42명 수령 (62.7%) — 최저 통과율", False),
        ("   • A6A 건물: 2명 전원 수령 (100%) — 최고 통과율", False),
    ]

    for i, (txt, is_bold) in enumerate(conclusions):
        if txt:
            add_text_box(slide, Inches(6.8), Inches(1.9) + Inches(i * 0.24), Inches(5.8), Inches(0.25),
                         txt, font_size=9, bold=is_bold,
                         color=NAVY if is_bold else DARK_GRAY)

    # Footer note
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
                 "※ 인센티브는 10개 조건 100% 충족 시에만 지급됩니다. 데이터 기준일: 2026-03-10",
                 font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✅ 보고서 생성 완료: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_report()
